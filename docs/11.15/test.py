from pptx import Presentation
from bs4 import BeautifulSoup
import mistune

# read md2 file 
def readMD(filename):
    md_text = open(filename, 'r').read()
    md = mistune.markdown(md_text)
    return BeautifulSoup(md, 'html.parser')

bs = readMD("md2ppt.md")

prs = Presentation()
default_font_size = 20

# Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = bs.h1.text
#subtitle.text = "hello, world"

def recursive_list(body_frame, contents, depth):

    def _addText(text, depth):
            p = body_frame.add_paragraph()
            p.text = text
            p.level = depth

    for li in contents.find_all('li', recursive=False):
        ul = li.find('ul')
        if li.find('ul'):
            _addText(next(li.stripped_strings), depth)
            recursive_list(body_frame, li.find('ul'), depth+1)
        elif li.find('ol'):
            _addText(next(li.stripped_strings), depth)
            recursive_list(body_frame, li.find('ol'), depth+1)
        else:
            _addText(li.text, depth)

slide_title = bs.h2
while slide_title: 
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]

    title.text = slide_title.text

    tf = body.text_frame

    slide_contents = slide_title.findNext()

    if slide_contents == slide_title.find_next('p'):
        path = str(slide_contents()).split(" ")[2].split('"')[1]
        pic = slide.shapes.add_picture(path, 21, 20)

    recursive_list(tf, slide_contents, 0)

    slide_title = slide_title.find_next('h2')

prs.save('out.pptx')
